#!/usr/bin/env python3
"""
Generate UW-branded poster (48" x 36" landscape) for brain_paper.
Layout documented in paper-writing/POSTER_LAYOUT.md.

Row 2 groups by experiment type (impulse | controller), not by analysis type.
"""

import os
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

PAPER     = r'C:\Users\aditya\Documents\projects\brain_paper\paper'
PDFTOPPM  = r'C:\Users\aditya\AppData\Local\Programs\MiKTeX\miktex\bin\x64\pdftoppm.exe'

# ── UW Brand Colours ──────────────────────────────────────────────────
PURPLE  = RGBColor(0x4B, 0x2E, 0x83)
GOLD    = RGBColor(0x85, 0x75, 0x4D)
GOLD_LT = RGBColor(0xD9, 0xCC, 0xAB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY= RGBColor(0xCC, 0xCC, 0xCC)
DARK    = RGBColor(0x22, 0x22, 0x22)

# ─────────────────────────────────────────────────────────────────────
# PDF → PNG CONVERSION (run before slide generation)
# ─────────────────────────────────────────────────────────────────────
def convert_pdf(pdf_name, dpi=220):
    """Convert first page of a PDF in paper/ to PNG. Returns filename or None."""
    pdf_path  = os.path.join(PAPER, pdf_name)
    base      = os.path.splitext(pdf_name)[0]
    prefix    = os.path.join(PAPER, base + '_poster')
    png_path  = prefix + '-1.png'
    if not os.path.exists(png_path):
        result = subprocess.run(
            [PDFTOPPM, '-r', str(dpi), '-png', '-singlefile', pdf_path, prefix],
            capture_output=True
        )
        # pdftoppm with -singlefile writes prefix.png (no -1 suffix)
        single = prefix + '.png'
        if os.path.exists(single):
            os.rename(single, png_path)
    if os.path.exists(png_path):
        print(f'  converted: {os.path.basename(png_path)}')
        return os.path.basename(png_path)
    print(f'  WARNING: conversion failed for {pdf_name}')
    return None

print('Converting PDFs...')
IMG = {}
IMG['imp_response_median']    = convert_pdf('imp_response_median.pdf')
IMG['tf_data_vs_model']       = convert_pdf('tf_data_vs_model_AL_0033_2025-01-29_en1.pdf')
IMG['tf_loao']                = convert_pdf('tf_loao_AL_0033_2025-01-29_en1.pdf')
IMG['ol_tf_three_sessions']   = convert_pdf('ol_tf_three_sessions.pdf')
# PNGs that exist directly
IMG['all_variance']           = 'all_variance_sessions.png'
IMG['all_MSE']                = 'all_MSE_sessions.png'
IMG['all_average']            = 'all_average_sessions.png'
IMG['trial_average_var']      = 'trial_average_var.png'
IMG['imp_motion_traces']      = 'imp_motion_traces.png'
IMG['motion_quartile']        = 'motion_quartile_combined.png'
IMG['freq_heatmap']           = 'freq_heatmap_combined.png'
IMG['imp_freqband']           = 'imp_freqband_AL_0033_2025-01-29_en1_amp5.png'

# ── Slide ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])
sh = slide.shapes

# ── Helpers ───────────────────────────────────────────────────────────
def box(x, y, w, h, fill=WHITE, line=MED_GRAY, lw=0.75):
    s = sh.add_shape(MSO.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(x, y, w, h, text, size=20, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def panel_hdr(x, y, w, title, hh=0.72, size=24):
    box(x, y, w, hh, fill=PURPLE, line=None)
    txt(x+0.15, y+0.04, w-0.3, hh-0.08, title,
        size=size, bold=True, color=WHITE)

def sub_hdr(x, y, w, title, size=19):
    txt(x, y, w, 0.48, title, size=size, bold=True, color=PURPLE)

def fig_place(x, y, w, h, img_key, caption):
    """Insert real image if available, else gray placeholder."""
    fname = IMG.get(img_key)
    if fname:
        fpath = os.path.join(PAPER, fname)
        if os.path.exists(fpath):
            sh.add_picture(fpath, Inches(x), Inches(y), Inches(w), Inches(h))
            return
    # fallback
    box(x, y, w, h, fill=LT_GRAY, line=MED_GRAY, lw=1.0)
    txt(x+0.1, y+h/2-0.35, w-0.2, 0.7, caption,
        size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

def fig_ph(x, y, w, h, caption):
    """Gray placeholder only (for items that need manual insertion)."""
    box(x, y, w, h, fill=LT_GRAY, line=MED_GRAY, lw=1.0)
    txt(x+0.1, y+h/2-0.35, w-0.2, 0.7, caption,
        size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

def body(x, y, w, h, text, size=17):
    txt(x, y, w, h, text, size=size, color=DARK, wrap=True)

def bullets(x, y, w, lines, size=17, gap=0.82):
    for line in lines:
        txt(x, y, w, gap+0.1, line, size=size, color=DARK, wrap=True)
        y += gap
    return y

# ─────────────────────────────────────────────────────────────────────
# WHITE SLIDE BACKGROUND
# ─────────────────────────────────────────────────────────────────────
box(0, 0, 48, 36, fill=WHITE, line=None)

# ─────────────────────────────────────────────────────────────────────
# TITLE BAR
# ─────────────────────────────────────────────────────────────────────
TH = 4.55
box(0, 0, 48, TH, fill=PURPLE, line=None)
box(0, TH, 48, 0.18, fill=GOLD, line=None)

for lx, lbl in [(0.28, "UW\nLOGO"), (43.9, "DEPT\nLOGO")]:
    box(lx, 0.28, 3.85, TH-0.56, fill=RGBColor(0x3D,0x24,0x6B), line=WHITE, lw=1.5)
    txt(lx, TH/2-0.5, 3.85, 1.0, lbl, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(4.45, 0.3, 39.1, 2.0,
    "Closed-Loop Widefield Control of Cortical Neural Activity via Optogenetics",
    size=49, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(4.45, 2.4, 39.1, 0.9,
    "A. Dewan¹  ·  N. Steinmetz¹       ¹Department of Neuroscience, University of Washington",
    size=26, color=GOLD_LT, align=PP_ALIGN.CENTER)
txt(4.45, 3.32, 39.1, 0.72,
    "adityad@uw.edu",
    size=20, color=GOLD_LT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# GRID CONSTANTS
# ─────────────────────────────────────────────────────────────────────
M  = 0.28
G  = 0.22

# Row 1
R1y = TH + 0.28
R1h = 19.4

C1x, C1w = M,               9.45
C2x, C2w = C1x+C1w+G,       18.0
C3x, C3w = C2x+C2w+G,       18.85

# Row 2  —  B1=Impulse brain state, B2=Controller brain state, B3=Conclusions
R2y = R1y + R1h + G
R2h = 36 - R2y - M - 0.33

B1x, B1w = M,               18.0
B2x, B2w = B1x+B1w+G,       18.0
B3x, B3w = B2x+B2w+G,       11.0

# ─────────────────────────────────────────────────────────────────────
# C1 — INTRODUCTION + SYSTEM
# ─────────────────────────────────────────────────────────────────────
box(C1x, R1y, C1w, R1h)
panel_hdr(C1x, R1y, C1w, "Introduction + System")

y = R1y + 0.85
sub_hdr(C1x+0.18, y, C1w-0.36, "Motivation"); y += 0.55
y = bullets(C1x+0.18, y, C1w-0.36, [
    "• Spontaneous cortical fluctuations drive trial-to-trial variability",
    "• Pre-stimulus brain state (motion, oscillations) predicts response amplitude",
    "• Open-loop stimulation cannot compensate for ongoing neural state",
], size=16, gap=0.78)

y += 0.2
sub_hdr(C1x+0.18, y, C1w-0.36, "Goal"); y += 0.55
body(C1x+0.18, y, C1w-0.36, 2.3,
     "Drive cortical ΔF/F to a fixed reference (−5%) using real-time widefield "
     "imaging feedback, reducing MSE and decoupling trial outcome from "
     "pre-stimulus brain state.", size=16); y += 2.5

sub_hdr(C1x+0.18, y, C1w-0.36, "PI Controller"); y += 0.55
y = bullets(C1x+0.18, y, C1w-0.36, [
    "• Discrete-time PI feedback controller",
    "• Feedback: kernel-avg ΔF/F in cortical ROI",
    "• Output: laser power (optogenetic inhibition)",
    "• Reference: −5% ΔF/F",
    "• Latency: ~80 ms  ·  Frame rate: 35 Hz",
], size=16, gap=0.72)

y += 0.18
sub_hdr(C1x+0.18, y, C1w-0.36, "Dataset"); y += 0.52
body(C1x+0.18, y, C1w-0.36, 1.0,
     "2 mice  ·  11 sessions\nAL_0033 (8)  ·  AL_0039 (3)", size=16); y += 1.2

# Experimental schematic — user will insert manually
fig_ph(C1x+0.18, y, C1w-0.36, R1y+R1h-y-0.18,
       "← Add experimental schematic here\n(widefield + laser setup)")

# ─────────────────────────────────────────────────────────────────────
# C2 — LTI VALIDATION
# ─────────────────────────────────────────────────────────────────────
box(C2x, R1y, C2w, R1h)
panel_hdr(C2x, R1y, C2w, "LTI Validation — Impulse + Step", size=24)

IX  = C2x + 0.18
IW  = C2w - 0.36
y   = R1y + 0.85

HALF    = (IW - G) / 2
LX      = IX
RX      = IX + HALF + G
TOP_H   = 12.1

# Left: dose-response
sub_hdr(LX, y, HALF, "Dose-Response Linearity")
body(LX, y+0.52, HALF, 1.35,
     "Peak inhibition energy (∫ΔF/F, 0–200 ms) scales linearly with amplitude "
     "(R²>0.97, n=3 sessions) — justifies LTI assumption.", size=15)
fig_place(LX, y+1.95, HALF, TOP_H-1.95,
          'imp_response_median', "paper/imp_response_median.pdf")

# Right: impulse TF (two panels)
sub_hdr(RX, y, HALF, "Impulse Transfer Function")
body(RX, y+0.52, HALF, 1.35,
     "Amplitude-normalised response fit with low-order LTI model "
     "(tfest, AIC-selected). LOAO cross-validation across amplitudes.", size=15)
TF_FIG_H = TOP_H - 1.95
TF_W     = (HALF - G) / 2
fig_place(RX,         y+1.95, TF_W, TF_FIG_H,
          'tf_data_vs_model', "tf_data_vs_model\n_AL_0033_en1.pdf")
fig_place(RX+TF_W+G,  y+1.95, TF_W, TF_FIG_H,
          'tf_loao', "tf_loao\n_AL_0033_en1.pdf")

# Bottom: step TF (full width)
sy = R1y + 0.72 + TOP_H + 0.3
box(IX-0.04, sy-0.08, IW+0.08, 0.04, fill=GOLD, line=None)
sy += 0.12

sub_hdr(IX, sy, IW, "OL Step Response — Transfer Function (3 Sessions)")
body(IX, sy+0.52, IW, 1.0,
     "Open-loop step response modelled with the same LTI plant; time constants "
     "match impulse TF, validating LTI model across stimulus types.", size=15)
fig_place(IX, sy+1.6, IW, R1y+R1h-sy-1.65,
          'ol_tf_three_sessions', "paper/ol_tf_three_sessions.pdf")

# ─────────────────────────────────────────────────────────────────────
# C3 — CONTROLLER PERFORMANCE
# ─────────────────────────────────────────────────────────────────────
box(C3x, R1y, C3w, R1h)
panel_hdr(C3x, R1y, C3w, "Controller Performance — CL vs OL (11 Sessions)", size=24)

IX3 = C3x + 0.18
IW3 = C3w - 0.36
MAIN_H = 12.9
FIG_W  = (IW3 - 2*G) / 3

y3 = R1y + 0.85

for i, (title, blurb, key, cap) in enumerate([
    ("Variance Suppression",
     "CL (green) suppresses trial-to-trial variance vs OL (red) — all sessions.",
     'all_variance', "all_variance_sessions.png"),
    ("MSE Reduction",
     "Per-session violin (t = +1 → +3 s): CL MSE < OL in every session.",
     'all_MSE', "all_MSE_sessions.png"),
    ("Trial-Average Convergence",
     "Mean trial trace: CL converges toward −5% reference; OL diverges.",
     'all_average', "all_average_sessions.png"),
]):
    cx = IX3 + i*(FIG_W+G)
    sub_hdr(cx, y3, FIG_W, title, size=17)
    body(cx, y3+0.5, FIG_W, 1.1, blurb, size=14)
    fig_place(cx, y3+1.65, FIG_W, MAIN_H-1.7, key, cap)

# Sine-wave strip
sy3 = R1y + 0.72 + MAIN_H + 0.3
box(IX3-0.04, sy3-0.08, IW3+0.08, 0.04, fill=GOLD, line=None)
sy3 += 0.12

sub_hdr(IX3, sy3, IW3, "Variable Reference Tracking — Sine Wave Controller", size=19)
body(IX3, sy3+0.52, IW3*0.38, 1.1,
     "Controller tracks a sinusoidal reference waveform. Trial-average ΔF/F "
     "follows target with low residual error (AL_0041 session).", size=15)
fig_place(IX3 + IW3*0.40, sy3+0.52, IW3*0.58, R1y+R1h-sy3-0.57,
          'trial_average_var', "trial_average_var.png")

# ─────────────────────────────────────────────────────────────────────
# ROW 2  —  B1 = Impulse brain state, B2 = Controller brain state
# ─────────────────────────────────────────────────────────────────────

BHALF = (B1w - 0.36 - G) / 2
y2 = R2y + 0.85

# ── B1: Impulse — Motion + Frequency ──────────────────────────────────
box(B1x, R2y, B1w, R2h)
panel_hdr(B1x, R2y, B1w, "Impulse — Brain State Dependence", size=23)

BLX1 = B1x + 0.18
BRX1 = BLX1 + BHALF + G

sub_hdr(BLX1, y2, BHALF, "Motion: Inhibition by Tertile", size=17)
body(BLX1, y2+0.5, BHALF, 1.2,
     "Low / mid / high motion trials: high motion reduces inhibition depth "
     "and increases inter-trial variability.", size=15)
fig_place(BLX1, y2+1.75, BHALF, R2y+R2h-y2-1.8,
          'imp_motion_traces', "imp_motion_traces.png")

sub_hdr(BRX1, y2, BHALF, "Frequency: Pre-Stim Spectral Power", size=17)
body(BRX1, y2+0.5, BHALF, 1.2,
     "Pre-stimulus 2–4 Hz power predicts inhibition depth. "
     "Representative mid-amplitude session (AL_0033).", size=15)
fig_place(BRX1, y2+1.75, BHALF, R2y+R2h-y2-1.8,
          'imp_freqband', "imp_freqband_..._amp5.png")

# ── B2: Controller — Motion + Frequency ───────────────────────────────
box(B2x, R2y, B2w, R2h)
panel_hdr(B2x, R2y, B2w, "Controller — Brain State Dependence", size=23)

BLX2 = B2x + 0.18
BRX2 = BLX2 + BHALF + G

sub_hdr(BLX2, y2, BHALF, "Motion: MSE by Quartile", size=17)
body(BLX2, y2+0.5, BHALF, 1.2,
     "High-motion trials have elevated MSE in OL; CL reduces this dependence — "
     "closed-loop partially compensates for motion.", size=15)
fig_place(BLX2, y2+1.75, BHALF, R2y+R2h-y2-1.8,
          'motion_quartile', "motion_quartile_combined.png")

sub_hdr(BRX2, y2, BHALF, "Frequency: Spectral Heatmap (MSE-Sorted)", size=17)
body(BRX2, y2+0.5, BHALF, 1.2,
     "Trials sorted by MSE: high-error OL trials show elevated 2–4 Hz pre-stim "
     "power. CL decouples MSE from spectral state.", size=15)
fig_place(BRX2, y2+1.75, BHALF, R2y+R2h-y2-1.8,
          'freq_heatmap', "freq_heatmap_combined.png")

# ── B3: Conclusions ───────────────────────────────────────────────────
box(B3x, R2y, B3w, R2h)
panel_hdr(B3x, R2y, B3w, "Conclusions", size=23)

y3 = R2y + 0.85
sub_hdr(B3x+0.18, y3, B3w-0.36, "Key Findings", size=18); y3 += 0.55
y3 = bullets(B3x+0.18, y3, B3w-0.36, [
    "✓  Cortical response is LTI: inhibition energy scales linearly with amplitude (R²>0.97)",
    "✓  Same LTI dynamics hold for step inputs — validates controller plant model",
    "✓  PI controller reduces MSE and variance in all 11 sessions",
    "✓  CL decouples trial outcome from motion and spectral brain state",
    "✓  Controller tracks variable (sine-wave) reference waveforms",
], size=15, gap=0.92)

y3 += 0.2
sub_hdr(B3x+0.18, y3, B3w-0.36, "Future Work", size=18); y3 += 0.55
y3 = bullets(B3x+0.18, y3, B3w-0.36, [
    "→  Three-layer prediction model (spont → OL → CL)",
    "→  MPC-optimal laser sequence benchmark",
    "→  Curto & Issa synced/desynced trial sorting",
], size=15, gap=0.82)

y3 += 0.2
box(B3x+0.18, y3, B3w-0.36, 0.06, fill=GOLD, line=None); y3 += 0.22
sub_hdr(B3x+0.18, y3, B3w-0.36, "Acknowledgments", size=16); y3 += 0.5
body(B3x+0.18, y3, B3w-0.36, 1.5,
     "Steinmetz lab, Dept of Neuroscience, University of Washington. "
     "Funding: [grant]. All experiments conducted under approved IACUC protocols.",
     size=14)

# ─────────────────────────────────────────────────────────────────────
# FOOTER
# ─────────────────────────────────────────────────────────────────────
box(0, 35.67, 48, 0.33, fill=PURPLE, line=None)
txt(0, 35.62, 48, 0.38,
    "Department of Neuroscience  ·  University of Washington  ·  Seattle, WA  98195",
    size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
out = r'C:\Users\aditya\Documents\projects\brain_paper\poster_template.pptx'
prs.save(out)
print(f'\nSaved: {out}')
