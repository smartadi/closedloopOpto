"""Build the NeuroAI Seattle 2026 talk deck. Structure + content only; no theme yet."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

Image.MAX_IMAGE_PIXELS = None
HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
ROOT = os.path.dirname(HERE)

W, H = 13.333, 7.5
BLACK = RGBColor(0x11, 0x11, 0x11)
GREY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xB0, 0x00, 0x00)
BODY = HEAD = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


def slide(title=None, sub=None):
    """Title top-left (may wrap to 2 lines); timing marker parked top-RIGHT so it can
    never collide with a wrapped title."""
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(10.4), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        r = tf.paragraphs[0].add_run()
        r.text = title
        r.font.size, r.font.bold, r.font.name = Pt(28), True, HEAD
        r.font.color.rgb = BLACK
    if sub:
        tb = s.shapes.add_textbox(Inches(11.1), Inches(0.38), Inches(1.65), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = sub
        r.font.size, r.font.name, r.font.italic = Pt(12), BODY, True
        r.font.color.rgb = GREY
    return s


def bullets(s, items, x, y, w, h, size=15, gap=7):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl, bold, col = (list(it) + [0, False, None])[:4]
        else:
            txt, lvl, bold, col = it, 0, False, None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size if lvl == 0 else size - 1)
        r.font.name = BODY
        r.font.bold = bool(bold)
        r.font.color.rgb = col or (BLACK if lvl == 0 else GREY)
    return tb


def text(s, txt, x, y, w, h, size=15, bold=False, col=None, italic=False, align=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        if align:
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), BODY, bold, italic
        r.font.color.rgb = col or BLACK
    return tb


def pic(s, name, x, y, boxw, boxh, cap=None, top=False):
    """Fit image in box preserving aspect. `top` anchors to the box top instead of
    centring, so a row of images shares a baseline. Returns (x, y, w, h) placed."""
    path = os.path.join(IMG, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = boxw, boxw / ar
    if h > boxh:
        h, w = boxh, boxh * ar
    px = x + (boxw - w) / 2
    py = y if top else y + (boxh - h) / 2
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))
    if cap:
        text(s, cap, x, py + h + 0.08, boxw, 0.3, size=11, col=GREY,
             align=PP_ALIGN.CENTER)
    return (px, py, w, h)


def caption(s, txt, x, y, w):
    text(s, txt, x, y, w, 0.3, size=11, col=GREY, align=PP_ALIGN.CENTER)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()


def movie(s, fname, poster, x, y, w, h):
    try:
        s.shapes.add_movie(os.path.join(ROOT, "presentation", fname),
                           Inches(x), Inches(y), Inches(w), Inches(h),
                           poster_frame_image=os.path.join(IMG, poster),
                           mime_type="video/mp4")
    except Exception as e:
        pic(s, poster, x, y, w, h)
        print("  ! movie embed failed for", fname, "->", e)


# ---------------------------------------------------------------- 0 TITLE
s = slide()
text(s, "Closed-loop optogenetic control of\nmesoscale cortical population activity",
     0.9, 2.2, W - 1.8, 2.0, size=36, bold=True)
text(s, "Aditya  ·  NeuroAI in Seattle 2026  ·  Allen Institute  ·  20–21 August",
     0.9, 4.5, W - 1.8, 0.5, size=17, col=GREY)
notes(s, """
15 min slot -> 12 min talk + 3 min Q. CONFIRM with organisers whether 15 includes Q.
Meeting theme: MULTI-AREA INTERACTIONS.
Organisers: Fairhall, Svoboda, Pereira-Obilinovic, Orsborn, Tim Kim.
TODO: author list + affiliations.
""")

# ---------------------------------------------------------------- 1 OPENING
s = slide("Precise measurement + active control makes the control toolkit available", "1:15")
bullets(s, [
    ("We can measure mesoscale cortical activity precisely and continuously — "
     "widefield, cortex-wide, 35 Hz", 0, True),
    ("…and act on it continuously — spatially targeted optogenetics", 0, True),
    ("…and identify the plant between them", 0, True),
    ("That combination is the precondition for a control system.", 0, True),
], 0.6, 1.45, 6.2, 2.0, size=15)

text(s, "What it unlocks", 0.6, 3.35, 6.2, 0.3, size=15, bold=True)
bullets(s, [
    "Brain state as the feedback signal — not noise to average away",
    "Optimization — the command is solved for, not hand-tuned",
    "Minimum-energy control = least-perturbative manipulation",
    "Model-driven control — designed, not tuned",
], 0.6, 3.72, 6.2, 1.9, size=14)

text(s, "ANALOGY — powered descent", 7.3, 1.45, 5.4, 0.3, size=14, bold=True)
bullets(s, [
    "Thrust is one-sided — you push, not pull",
    "The laser only inhibits:  0 ≤ u ≤ u_max",
    "Minimum-fuel, not just 'any landing'",
    "Bounded non-convex thrust → convex program",
    ("⇒ constrained QP, not LQR", 0, True),
], 7.3, 1.8, 5.4, 2.1, size=13)

text(s, "…but with so much variability it's too hard.\nSo how do we figure it out?",
     7.3, 3.55, 5.4, 0.6, size=14, bold=True, italic=True)

text(s, "What already exists — and the gap", 7.3, 4.3, 5.4, 0.3, size=14, bold=True)
text(s, "Closed-loop optogenetics is not new: Newman 2015 (the optoclamp — also a PI "
        "controller, our direct ancestor), Bolus 2018/2021 (model-based, then LQR in "
        "awake mice), Kanta 2019, Clancy 2014, Grosenick 2015.",
     7.3, 4.68, 5.4, 1.0, size=12)
bullets(s, [
    ("One thread has the LOOP but not the SCALE — single-neuron or LFP, ephys "
     "feedback", 0),
    ("The other has the THEORY but no ACTUATOR — controllability on structural "
     "connectomes, inputs hypothetical", 0),
], 7.3, 5.47, 5.4, 1.2, size=12)

text(s, "Closing the loop turns the perturbation into a CONTROLLED INDEPENDENT "
        "VARIABLE — not “we delivered 2 mW to this region”, but “we held this region "
        "at −5 % ΔF/F”. Widefield + targeted opto gives loop, scale and theory at once.",
     0.6, 5.25, 6.2, 1.1, size=14, bold=True)
text(s, "The talk: identify the plant → close the loop → characterise the residual → "
        "model the disturbance.",
     0.6, 6.55, 12.1, 0.35, size=13, italic=True)
notes(s, """
OPPORTUNITY FRAMING, NOT COMPLAINT (reframed at user request).

The old deficiency point survives as ONE CLAUSE: open-loop delivers a fixed light
pattern irrespective of ongoing state, so the same command lands differently depending
on when it arrives (Harris & Thiele 2011) -- then IMMEDIATELY turn positive: that state
is measurable now, so it becomes the feedback signal rather than the noise.

MINIMUM ENERGY -- make the SCIENTIFIC argument, not the aesthetic one:
minimum energy = least-perturbative causal manipulation. Less light, less heating,
less off-target drive. Not "did the perturbation work" but "what is the SMALLEST push
that produces the effect?" That reframes optogenetics as a DESIGNED intervention.

ROCKET ANALOGY -- say it ONCE, crisply, then drop it. Half this room is
neuroscientists; an extended aerospace metaphor reads as engineering-imperialism.
The line: "one-sided bounded actuation with a fuel cost is a solved problem in
aerospace, and it is exactly the structure of a laser that can only inhibit."

DISANALOGIES (say these -- they are the insurance, and the interesting part):
 1. No plant for free. No first-principles model; identified from data, per animal.
 2. THE DISTURBANCE IS THE SYSTEM. Wind is exogenous to an aircraft; ongoing cortical
    activity is the thing you are studying. You reject a signal you also want to
    measure. <- this sentence turns the analogy into the talk's thesis.
 3. You don't know the reference. A rocket's trajectory comes from the mission. What
    trajectory SHOULD neural activity follow? Open scientific question.

CITATION GAP: connectome-controllability literature (Gu et al. 2015;
Pasqualetti/Bassett) is NOT in refs.bib -- verified. Do not cite from memory.

--- FOLDED IN FROM THE OLD GAP SLIDE (2026-08-18) ---
ACKNOWLEDGE GENEROUSLY AND SPECIFICALLY BEFORE CLAIMING THE GAP.
DO NOT FRAME OPEN-LOOP AS FLAWED. It produced the foundational results, and SVOBODA IS
AN ORGANISER AND IN THE ROOM. Frame as: right tool, different question.
Newman 2015 is worth naming aloud as the direct ancestor -- it is also a PI controller.
That generosity buys credibility for the gap claim that follows.
Source: Closedloop_edit/introduction.tex paragraphs 4-7.
The spine line ("once you are actively rejecting a disturbance, the disturbance is the
rest of the brain") now lives on the contra->ipsi map slide -- say it there, and again
in the close.
""")

# ---------------------------------------------------------------- 3 THE LOOP
s = slide("The loop", "0:45")
pic(s, "wfpath.png", 0.7, 1.5, 3.0, 4.3, cap="widefield light path", top=True)
pic(s, "svd_frame.png", 4.1, 1.5, 3.0, 4.3, cap="SVD readout frame", top=True)
pic(s, "arch.png", 7.4, 1.5, 5.3, 2.7, cap="control architecture", top=True)
text(s, "Widefield SVD readout @ 35 Hz  →  kernel-mean ΔF/F at a target pixel\n"
        "→  PI controller  →  638 / 594 nm laser",
     7.4, 4.5, 5.3, 0.9, size=15)
text(s, "~14 ms minimum loop latency\nReference = −5 % ΔF/F",
     7.4, 5.5, 5.3, 0.8, size=15, bold=True)
notes(s, """
0:45 ONLY -- slide 1 now carries the framing, this is just the implementation.
Trimmed from 1:00 to buy slide 1 its time.

Numbers: 35 Hz control rate, ~14 ms minimum loop latency, ref = -5 %dF/F.
Plant lag ~160 ms -- comes up again on slide 6, where the 143 ms preview lookahead is
matched to it BY DESIGN.
""")

# ---------------------------------------------------------------- 4 PLANT ID
s = slide("The plant is approximately linear — and it transfers", "1:15")
pic(s, "imp_single.png", 0.6, 1.45, 3.9, 2.9, top=True)
pic(s, "imp_response.png", 4.7, 1.45, 3.9, 2.9, top=True)
pic(s, "tf_shape.png", 8.8, 1.45, 3.9, 2.9, top=True)
bullets(s, [
    ("Input–output characterization per session reveals that LTI systems explain "
     "trial-averaged dynamics", 0),
    ("TF is roughly consistent across sessions", 0),
    ("similar time scales", 1),
    ("similar shapes — not strictly the same though.", 1),
], 0.6, 4.85, 12.1, 1.8, size=15)
notes(s, """
DO NOT say "LTI" as a proven waveform-superposition claim -- waveform superimposition
was NEVER verified. Say "approximately linear, well described by a low-order model".
n = 3 sessions for the dose-response.

TF-D (cross-session model-swap) uses FREE GAIN -- it is a claim about DYNAMICS, not
gain; gain is re-tuned per session by construction. Say so if you show it.

TF-D IS BACKUP-BY-DEFAULT -- dropped to buy slide 9 its time. tau-consistency (TF-A)
is the load-bearing half. Add TF-D back only if the slot turns out to be 15 + Q.

SIGNPOST FORWARD: "hold onto this -- MPC needs a forward model, and this is it."
""")

# ------------------- 4b CONTROLLER DESIGN
s = slide("The controller: feedforward inversion, then PI on top", "0:45")

text(s, "Feedforward", 0.6, 1.4, 6.3, 0.3, size=15, bold=True)
text(s, "u = Kr · e,      e = y_ref − y,      Kr = U Y†",
     0.6, 1.75, 6.3, 0.4, size=17, bold=True)
text(s, "Kr = least squares on the impulse calibration.", 0.6, 2.2, 6.3, 0.3, size=13)

text(s, "PI feedback", 0.6, 2.8, 6.3, 0.3, size=15, bold=True)
text(s, "u = Kr·e  +  Kp·e  +  Ki·Σ e Δt", 0.6, 3.15, 6.3, 0.4, size=17, bold=True)
text(s, "Open loop = the same law with Kp = Ki = 0. Feedback added to a calibrated "
        "feedforward — not feedback vs nothing.",
     0.6, 3.6, 6.3, 0.7, size=14, bold=True)

bullets(s, [
    ("Command clipped 0–5 V (≈ 36 mW) — the one-sided bound, in hardware", 0),
    ("Anti-windup clamp on the accumulator", 0),
    ("s = [x ; ξ],  ξ = Σ e Δt,  u = K s — the integrator is a STATE", 0, True),
], 0.6, 4.5, 6.3, 1.5, size=13)

pic(s, "tune_grid.png", 7.2, 1.5, 2.9, 2.7, top=True)
pic(s, "tune_auto_cost.png", 10.2, 1.5, 2.6, 2.5, top=True)
caption(s, "cost surface J(Kp,Ki), AL_0033 03-05", 7.2, 4.3, 2.9)
caption(s, "online auto-tune, same mouse — 16.6 → 12.3", 10.2, 4.3, 2.6)

text(s, "Gains tuned per session, model-free", 7.2, 5.0, 5.5, 0.3, size=15, bold=True)
text(s, "Grid basin ≈ (0.05, 0.10); online search settles at (0.068, 0.064).",
     7.2, 5.35, 5.5, 0.6, size=14)
text(s, "No model needed to tune PI. A model is what it takes to beat it.",
     7.2, 6.0, 5.5, 0.6, size=14, bold=True, italic=True)
caption(s, "J = mean ‖y − y_ref‖₂, t = 0–3 s, ref = −5 % ΔF/F", 7.2, 6.7, 5.5)

notes(s, """
ADDED 2026-08-18 at user request -- "a slide is missing between 4 and 5 showing the
actual controller design". Equations are Methods eq:ff / eq:kr / eq:pi / eq:aug_state
from Closedloop_edit/methods_edit.tex (sec:ff_gain, sec:pi_controller_design).

THE ONE SENTENCE THAT MUST NOT BE SKIPPED: open loop is the SAME feedforward law with
Kp = Ki = 0. If the room thinks OL = "no controller / no light", every OL-vs-CL number
later in the talk looks inflated. Say it once, clearly, here.

WHY THIS SLIDE EARNS 0:45: it is where the aerospace framing cashes out in hardware --
0 <= u <= u_max is literally the 0-5 V clip -- and it plants the integrator-as-a-state
that the internal model principle slide pays off later. Do not just read the equations.

THE PANELS ARE METHODS-ONLY (PAPER.md fig:cost_landscape, no Results figure number).
 - Grid: gain_grid.m, AL_0033 2026-03-05. Clean interior min near (0.05, 0.10).
 - Auto-tune: greedy zeroth-order, accept-if-cost-lowered, AL_0033 2026-03-17,
   accepted gains only (Kdata/Kval -- NOT input_params, which logs rejected probes).
   Cost 16.6 -> 12.3, settles at (0.068, 0.064).
 - BOTH ARE AL_0033 ON PURPOSE. Grid-vs-autotune comparisons are SAME-MOUSE ONLY, and
   showing only AL_0033 also avoids introducing AL_0034, which is a tuning-only mouse
   that appears nowhere else in the talk.
 - Cost = mean ||y - ref||_2 over t = 0-3 s, y = online kernel-mean %dF/F, ref = -5.

IF ASKED "why not just LQR / MPC from the start?": there is a backup slide. Short
answer -- we did not have a per-animal forward model when the loop was built; PI needs
no model, and the grid+autotune agreement shows it can be tuned honestly without one.""")

# ---------------------------------------------------------------- 5 CL VIDEO
s = slide("Closed loop holds the set-point", "1:45  ·  VIDEO")
movie(s, "controller_demo_m13_full_hd.mp4", "poster_ctrl.png", 0.6, 1.6, 7.6, 4.28)
caption(s, "controller_demo_m13_full_hd.mp4  ·  0:39  ·  m13 = AL_0039 2025-04-20, "
           "step, 100 trials", 0.6, 5.95, 7.6)
pic(s, "pooled_rmse.png", 8.7, 1.7, 4.0, 3.5)
caption(s, "3K  per-session median OL vs CL RMSE, 15 sessions", 8.7, 5.3, 4.0)
text(s, "CL < OL in 14 of 15 sessions\nsignrank p = 1.2 × 10⁻⁴",
     8.7, 5.75, 4.0, 0.8, size=16, bold=True, align=PP_ALIGN.CENTER)
notes(s, """
THE VIDEO REPLACES static panels 3A/3B/3C -- it shows single trial, trial-average and
evolving across-trial variance AT ONCE and IN MOTION. Keep 3A-3C as an AV-failure
fallback slide.

SAY THE HONESTY LINE ONCE: "this is a replay of real session data through a replica of
the live GUI, not a screen capture." Traces, gains, RMSE and the widefield movie are
all real session data; the layout is restyled for presentation.

NUMBERS: CL < OL in 14/15 sessions, signrank p = 1.2e-4. Pooled OL 3.11 / CL 2.44
%dF/F. RMSE window t = 0 -> +3 s. m13 has a ~30% CL-vs-OL gap.

SAY RMSE, NEVER MSE (project-wide rename 2026-07-16).

If asked why the OL command looks like a clean step: the laser panel shows the command
AMPLITUDE ENVELOPE (sliding-window max), not the raw carrier.
""")

# ------------------- 6 WHAT DRIVES THE TRACKING ERROR (FIG 4 block 2)
s = slide("What drives the tracking error — and when", "1:15")
pic(s, "f4_decomp.png", 0.5, 1.45, 4.1, 2.9,
    cap="unique (partial) R² of trial RMSE", top=True)
pic(s, "f4_initdev.png", 4.8, 1.45, 3.9, 2.9,
    cap="initial-deviation quartile", top=True)
pic(s, "f4_hi24.png", 8.9, 1.45, 3.9, 2.9,
    cap="relative 2–4 Hz quartile, settled window", top=True)
caption(s, "613 closed-loop trials / 11 sessions; 2000-bootstrap 95 % CI",
        0.5, 4.5, 12.3)

bullets(s, [
    ("Early (0–1 s) the error is the TRANSIENT — where the brain happened to be at "
     "onset. Unique R² 0.38", 0, True),
    ("Late (1–3 s) that collapses 13-fold to 0.03. The controller settles the "
     "transient out", 0),
    ("Motion contributes essentially nothing — unique R² 0.001 in both windows", 0),
    ("What is left late is the RELATIVE 2–4 Hz state: 0.09 → 0.12, which is 93 % of "
     "the explained variance by the settled window", 0, True),
], 0.5, 4.9, 7.6, 2.2, size=14)

text(s, "So the residual error is not a tuning\nproblem and it is not movement.",
     8.5, 4.9, 4.3, 0.8, size=15, bold=True)
text(s, "It is one specific band of ongoing cortical activity — which is exactly what "
        "the next slide is about.",
     8.5, 5.8, 4.3, 1.0, size=14, italic=True)
notes(s, """
FIG-4 PART 1, the error-contribution model. Panels from paper/images/figure4/:
f4p1_error_decomp.png (cl_factor_decomp_panel.m), claim3_initdev_early_late.pdf,
claim2_delta_hi24_late.pdf.

NUMBERS -- unique (partial) R^2, 613 CL trials / 11 sessions, 2000-boot 95% CI:
   initial deviation  0.381 [0.299 0.458] EARLY 0-1 s -> 0.029 [0.004 0.077] LATE 1-3 s
   motion             0.001 -> 0.001
   relative 2-4 Hz    0.089 [0.054 0.130] -> 0.122 [0.077 0.176]
   full model R^2     0.407 early / 0.132 late
   ==> relative 2-4 Hz is 22% of explained variance early, 93% LATE.

Init-dev quartile panel: transient 2.40 -> 4.13 across Q1->Q4; settled 1.95 -> 2.44.
The transient window is strongly graded by where the trial started; the settled window
is nearly flat. That IS the 13x collapse, made visual.

ROBUSTNESS -- the motion bar is a real null, not missing data: m1/m3/m7/m8 have no
motion recording, so gating the pool on motion drops a third of the data. Re-scored
with the gate OFF (852 trials / 15 sessions) full R^2 = 0.389/0.123 vs 0.407/0.132
gated -- nothing moves. That is what licenses saying "motion does not matter" rather
than "motion was not measured here". Say it if challenged.

⚠ PANEL HYGIENE: claim2_delta_hi24_late is safe here as the 2-4 Hz MAGNITUDE
companion, on its own. NEVER place claim2_delta_lo12_late beside it -- in the pooled-
quartile form the 1-2 Hz series rises almost identically (between-session baseline
structure) and the two together argue AGAINST the dissociation. The 1-2 Hz null lives
on the NEXT slide, in paired-slope form, where picture and statistic agree.

THIS SLIDE SETS UP THE WATERBED. It establishes WHAT is left in the error. The next
slide says what the loop does to it -- and that it makes it worse.
""")

# ------------------- 7 WHAT THE LOOP CAN AND CANNOT REJECT (FIG 4 block 2)
s = slide("What the loop can reject — and what it makes worse", "1:30")
pic(s, "f4_slope_motion.png", 0.5, 1.4, 3.5, 2.5,
    cap="motion  OL +0.69 → CL −0.02   p = 0.002", top=True)
pic(s, "f4_slope_delta24.png", 4.2, 1.4, 3.5, 2.5,
    cap="relative 2–4 Hz  OL −0.51 → CL +0.81   p = 0.013", top=True)
pic(s, "f4_slope_delta12.png", 7.9, 1.4, 3.5, 2.5,
    cap="relative 1–2 Hz  null control   p = 0.76 n.s.", top=True)
caption(s, "per-session slope of trial RMSE on each state factor, OL vs CL — "
           "15 sessions, four mice (motion n = 11)", 0.5, 4.35, 10.9)

text(s, "Internal model principle: to reject a disturbance, the loop must contain a model of its dynamics.", 0.5, 4.85, 6.9, 0.55,
     size=15, bold=True)
bullets(s, [
    ("Motion is slow and quasi-DC. The integrator IS a model of a constant — so "
     "feedback decouples it completely (+0.69 → −0.02)", 0),
    ("The loop has NO oscillatory mode at 2–4 Hz. The slope does not just survive, "
     "it FLIPS SIGN — feedback makes that band WORSE", 0, True),
    ("1–2 Hz is untouched, so this is band-specific, not generic", 0),
], 0.5, 5.45, 6.9, 1.8, size=13)

text(s, "That sign flip is the waterbed", 7.7, 4.9, 5.1, 0.3, size=15, bold=True)
text(s, "A controller merely lacking a model of that band would leave the slope alone "
        "(S ≈ 1: no rejection, no harm). A flip means the loop is paying for its "
        "low-frequency suppression — Bode's integral, not a null result.",
     7.7, 5.3, 5.1, 1.4, size=13)
text(s, "Unique R² of trial error: initial deviation 0.38 → 0.03 late, motion 0.001, "
        "relative 2–4 Hz 0.09 → 0.12 — by late trial it is 93 % of the explained "
        "variance.",
     7.7, 6.5, 5.1, 0.9, size=13, bold=True)
notes(s, """
FIG-4 PART 1 (PAPER.md "Block 2 / panel 4B"). Added 2026-08-18 at user request.
This is the IMP argument WITH DATA, and it is the bridge from slide 10 to the close.

THE ARGUMENT (RESEARCH 2026-08-13):
 - MOTION is a slow, quasi-DC nuisance. A PI integrator IS an internal model of a
   constant, so the loop should reject it -- and it does, completely:
   OL +0.686 -> CL -0.016, signrank p = 0.0020 **.
 - RELATIVE 2-4 Hz: the loop contains NO oscillatory mode at that frequency. The
   key point is that the slope does not merely SURVIVE, it FLIPS SIGN:
   OL -0.511 -> CL +0.813, p = 0.0125 *. A controller that simply lacked a model of
   that band would leave the OL slope untouched (S ~ 1: no rejection, no harm).
   A flip means feedback is ACTIVELY MAKING 2-4 Hz WORSE. That is the Bode sensitivity
   integral / waterbed effect -- suppression at low frequency is paid for higher up.
 - 1-2 Hz NULL CONTROL: OL -0.319 -> CL +0.068, p = 0.76 n.s. The dissociation is
   band-specific. n=15 sessions, four mice.

ERROR DECOMPOSITION (613 CL trials / 11 sessions, 2000-boot 95% CI; unique partial R2):
   initial deviation 0.381 [0.299 0.458] EARLY -> 0.029 [0.004 0.077] LATE (13x collapse)
   motion            0.001 -> 0.001   (nothing)
   relative 2-4 Hz   0.089 [0.054 0.130] -> 0.122 [0.077 0.176]
   full model R2 0.407 early / 0.132 late.
   Relative 2-4 Hz = 22% of explained variance early, 93% LATE.
 ROBUSTNESS: with the motion gate off (852 trials / 15 sessions) R2 = 0.389/0.123 vs
 0.407/0.132 -- nothing moves. That is what licenses reading the motion bar as "motion
 does not matter" rather than "motion was not measured here".

⚠ PANEL HYGIENE:
 - The 1-2 Hz null MUST be shown as `factor_slope_delta12` (per-session paired slope),
   NOT as `claim2_delta_lo12_late`. In the pooled-quartile form the 1-2 Hz series rises
   almost identically to 2-4 Hz (between-session baseline structure), so side by side
   those two argue AGAINST the dissociation. Only the paired-slope form has picture and
   statistic agreeing (p=0.0125 vs p=0.76).
 - Known cosmetic issue: factor_slope_motion and factor_slope_delta24 were exported with
   different y-ranges, which blunts the contrast between the two strongest panels. Not
   yet fixed at source.

WHY THIS EARNS THE CLOSE: it says exactly what is missing. The loop models DC and
rejects it; it has no model of the oscillatory band and pays for that. Give it a model
of the network's dynamics -- that is MPC.
""")

# ------------------- 8 SETUP: THE CONTRA→IPSI MAP (FIG 4 block 3)
s = slide("How we build the disturbance: a stim-blind contra → ipsi map", "0:45")
pic(s, "kernel_paint.png", 0.5, 1.35, 6.1, 5.6, top=True)
caption(s, "six ipsi target sites (right, ringed) and their contra kernels — "
           "opacity ∝ |weight|", 0.5, 7.0, 6.1)

text(s, "Reading the map", 7.1, 1.4, 5.6, 0.3, size=16, bold=True)
text(s, "Each ipsi site is predicted by a distinct, spatially coherent territory of "
        "contralateral cortex. The dots are the sparse predictor grid, coloured by "
        "which ipsi target that contra pixel best predicts.",
     7.1, 1.8, 5.6, 1.3, size=15)
text(s, "The mapping is topographic — not one diffuse global signal.",
     7.1, 3.15, 5.6, 0.5, size=16, bold=True)

text(s, "What we do", 7.1, 4.05, 5.6, 0.3, size=16, bold=True)
bullets(s, [
    "Fit contra → ipsi on SPONTANEOUS, laser-off frames only",
    "Use only contra pixels the laser does not affect ⇒ stim-blind by construction",
    "Held-out contra → ipsi R² = 0.87 (m4); 0.978 on Ye et al. 2023's own data",
    "Deploy during trials → the prediction is Global, the counterfactual ipsi "
    "“as if no laser”.   Local = Actual − Global",
], 7.1, 4.45, 5.6, 2.2, size=15)
text(s, "Once you are actively rejecting a disturbance, the disturbance is the rest "
        "of the brain.",
     7.1, 6.6, 5.6, 0.8, size=15, bold=True, italic=True)
notes(s, """
SETUP SLIDE -- added 2026-08-18 at user request, placed before the rejection slide
because Global has to be defined before ER can mean anything.

WHAT TO SAY:
 - The predictor is fit on SPONTANEOUS, laser-off frames, using only contra pixels that
   the laser does not affect. So it is stim-blind BY CONSTRUCTION, not by correction.
 - Deployed during a trial it produces GLOBAL: what the ipsi site would have done with
   no laser. LOCAL = Actual - Global is then the controller's own contribution.
 - Held-out contra->ipsi R^2 = 0.874 on m4 (rho = 0.936). Validated against Ye et al.
   2023 -- 0.978 on their data, reproducing their ~0.99.
 - The left map answers "which contra pixels carry information about this ipsi site":
   structured, posterior/lateral, not uniform. The + is the target site.

THE PAINT MAP -- source: impulse-analysis/kernelmap_paint.m -> paper/kernelmap_paint.png
(2026-07-10). Zhiwen-style. 6 ipsi targets laid out in a 1-2-3 triangle, Okabe-Ito
colourblind-safe palette. Each contra kernel is airbrush-PAINTED with opacity = |weight|
(Gaussian-smoothed); the sparse predictor grid sits on top, each grid pixel coloured
winner-take-all by which ipsi target it most strongly predicts. Built on the IMPULSE
sparse-OLS predictor (ols_pixel_predictor_wip.m), not the controller Stage-1/2 fit --
same method, different pipeline. Say "one session" if pressed.

THE POINT TO LAND: the contra->ipsi mapping is TOPOGRAPHIC. Different ipsi sites draw
from different contra territory, in spatially coherent blocks. That matters twice over:
it says Global is structured rather than one diffuse global signal, and it is itself a
multi-area interaction result -- which is this meeting's theme.

BACKUP, IF ASKED "are those weights really where the signal comes from?":
 A regression FILTER is optimised to PREDICT, not to LOCALISE -- it will put large
 opposite-sign weights on correlated neighbours to cancel noise. The Haufe forward
 pattern a = cov(x, yhat) is the interpretable object. On m4:
 rho(filter, pattern) = -0.020, 47% of 249 px flip sign; the filter map is
 salt-and-pepper, the pattern map a smooth posterior/lateral contra region.
 (Haufe et al. 2014, forward vs backward models.) The paint map's winner-take-all
 assignment is more robust than per-pixel weight signs, but it IS filter-derived --
 don't oversell individual pixels. Panels for this are in talk/img/kernel_filter.png
 and kernel_pattern.png if you want a backup slide.
""")

# ------------------- 9 REJECTION + INTERNAL MODEL (FIG 4 block 3)
s = slide("Disturbance rejection — and what the internal model principle says next",
          "1:30")
pic(s, "rej_demo_ol.png", 0.6, 1.45, 3.3, 2.0, cap="OPEN LOOP", top=True)
pic(s, "rej_demo_cl.png", 0.6, 3.75, 3.3, 2.0, cap="CLOSED LOOP", top=True)
pic(s, "rej_paired.png", 4.2, 1.9, 3.1, 2.6, top=True)
text(s, "energy ratio  ER = ‖A − ref‖² / ‖G − ref‖²\n"
        "1 = no work done    < 1 = controller gain",
     4.2, 4.7, 3.4, 0.7, size=13, align=PP_ALIGN.CENTER)
text(s, "OL 0.291  →  CL 0.203\nbetter in 13 of 13 sessions\nsignrank p = 2.4 × 10⁻⁴",
     4.2, 5.5, 3.4, 1.0, size=15, bold=True, align=PP_ALIGN.CENTER)

bullets(s, [
    "Per-trial transmission β (slope of Actual on Global): m4 OL 0.97 → CL 0.56 — "
    "open-loop the disturbance passes through untouched; feedback nearly halves it",
    "Rejection is state-INVARIANT: a constant OL−CL offset across motion and δ "
    "quartiles",
], 7.9, 1.45, 4.8, 2.0, size=13)
text(s, "Internal model principle", 7.9, 3.5, 4.8, 0.3, size=15, bold=True)
text(s, "To reject a disturbance, the controller must contain a model of that "
        "disturbance's dynamics.",
     7.9, 3.85, 4.8, 0.8, size=14)
text(s, "A PI controller's integrator is an internal model of a CONSTANT. That is why "
        "it holds a set-point against DC drift — and why it stops at ER ≈ 0.20 against "
        "ongoing cortical activity, which is not constant.",
     7.9, 4.65, 4.8, 1.3, size=14)
text(s, "The residual is not a tuning failure.\nIt is a missing model.",
     7.9, 6.0, 4.8, 0.7, size=15, bold=True, italic=True)
caption(s, "panels are UNGATED (no R² floor, n = 13)", 0.6, 6.3, 3.3)
notes(s, """
THE EMPIRICAL CENTRE OF THE SECOND HALF -- and the strongest cross-session result in
the talk. Promoted 2026-08-18 at user request.

THE METRIC: ER = ||A-ref||^2 / ||G-ref||^2, Nick's 2026-07-28 energy ratio, implemented
and PRIMARY. Both terms share the reference, so ER = 1 means the controller did no work
and the LEVEL is interpretable on its own -- not just the OL-vs-CL contrast.
The legacy rho = ||A-ref||/||G|| is retained for reproducibility. NEVER MIX THE TWO --
every number logged before 2026-08-10 is in rho units.

NUMBERS (all run, all cross-session):
 - ER median OL 0.291 -> CL 0.203; CL better 13/13; signrank p = 2.4e-4.
 - Per-session rank-sum significant in 11/13. The two misses are exactly the two new
   mice with weak predictors (AL_0048 p=0.299 R2=0.441; AL_0051 p=0.755 R2=-2.058), so
   the result does NOT rest on them -- their inclusion is CONSERVATIVE. Say that.
 - Strongest: AL_0033_0304_e1 0.477->0.207 p=1.6e-7; AL_0039_0430_e3 0.311->0.182
   p=6.1e-8.
 - Transmission beta m4: OL 0.97 -> CL 0.56.

STATE-INVARIANCE IS THE HONEST FINDING, and it is a good one. State-dependence returned
NULL: [IMP-STATE-QUARTILE] on m4 gave all four r_s null; [IMP-XSTATE] across 7 sessions
found no surviving state slope, Friedman n.s. Do not imply a state effect.

CAVEAT ON THE PANELS: they carry a burned-in "UNGATED (no R^2 floor, n=13)" annotation
that collides with the panel title, so the deck crops the top band -- THE CAVEAT IS IN
THE SLIDE CAPTION INSTEAD. Do not drop it.

THE PAYLOAD IS THE THEORY. Internal model principle (Francis & Wonham): to reject a
class of disturbances the controller must contain a model of their dynamics. PI's
integrator = internal model of a CONSTANT. Ongoing cortical activity is not constant.
So the floor at ER ~ 0.20 is not bad tuning -- it is a MISSING MODEL.
==> MPC is not an add-on. It is what IMP prescribes. This sets up the close.
""")

# ------------------- 10 SINE / PREVIEW VIDEO (FIG 5)
s = slide("Tracking a moving reference: feedback vs preview", "1:15  ·  VIDEO")
movie(s, "ff_demo_0721_hd.mp4", "poster_ff.png", 0.6, 1.6, 6.6, 3.71)
caption(s, "ff_demo_0721_hd.mp4  ·  TRIM TO ~30 s", 0.6, 5.4, 6.6)
pic(s, "sine_rmse.png", 7.6, 1.6, 2.4, 3.1, cap="5I  total RMSE", top=True)
pic(s, "sine_phase.png", 10.3, 1.6, 2.4, 3.1, cap="5J  1 Hz phase lag", top=True)
text(s, "The error win comes from FEEDBACK.\nThe phase win comes from PREVIEW.",
     7.6, 5.15, 5.1, 0.7, size=16, bold=True)
text(s, "CL < OL in all 3 sessions, pooled p = 1.5 × 10⁻⁴.   Preview: 60° → 5° lag; "
        "143 ms lookahead ≈ the measured ~160 ms plant lag.",
     0.6, 6.1, 12.1, 0.5, size=13)
text(s, "Preliminary: n = 3 sessions, one mouse, one hemisphere.",
     0.6, 6.65, 12.1, 0.3, size=13, col=RED, italic=True)
notes(s, """
DO NOT SAY:
 - that preview reduces error. OL vs OL+preview is N.S., p = 0.49 (p = 0.32 on s2).
 - that CL+preview beats CL. p = 0.093 IN CL's FAVOUR.
 - that Friedman across sessions is significant. n = 3, p = 0.060, UNDERPOWERED.
 - "significant in every session" -- s3 alone is n.s. (p = 0.36).

The supportable sentence, if you want CL+preview: "CL+preview is the only mode that is
both significantly below open-loop on error AND lag-free."

CAVEAT ALOUD: one mouse (AL_0048), one hemisphere (right/inhibitory), 3 sessions.
Nick has asked for a water-restricted repeat before any stats claim. Call it
preliminary -- do not let the room think this is settled.

WHY THIS SLIDE STILL EARNS ITS PLACE NOW THAT THE GRID SLIDE IS GONE: it is the only
place in the talk where LOOKAHEAD is shown to work on real data. Preview at a 143 ms
horizon cancels the plant's phase lag -- an existence proof for anticipation on a KNOWN
reference, which the close then generalises to a PREDICTED disturbance. Signpost it.
""")

# ------------------- 11 CLOSE
s = slide("Model-predictive control does better — and needs less light", "1:15")
pic(s, "opt_tracking.png", 0.6, 1.45, 4.0, 3.3, top=True)
pic(s, "opt_command.png", 4.8, 1.45, 3.9, 3.3, top=True)
caption(s, "m4, trial-averaged. PI overshoots to −6.5 % and oscillates; "
           "the optimal controller sits on the reference — using less command.",
        0.6, 4.85, 8.1)

text(s, "Steady RMSE   PI 1.20  →  optimal 0.15", 0.6, 5.35, 8.1, 0.35,
     size=17, bold=True)
bullets(s, [
    "Mean laser command PI 1.563 → 1.300 — better tracking with LESS light. "
    "Least-perturbative and best-performing are not in tension here",
    "Causal 0.146 ≈ clairvoyant 0.000 — nearly all the gain is reachable CAUSALLY. "
    "The PI gap is not about seeing the future",
], 0.6, 5.8, 8.1, 1.3, size=13)

text(s, "The internal model principle told us this", 9.0, 1.45, 3.7, 0.3,
     size=15, bold=True)
text(s, "PI rejects a constant. Give the controller a model of the network's "
        "dynamics and it rejects what the network actually does.",
     9.0, 1.8, 3.7, 1.1, size=14)
text(s, "That model is where learned methods enter — and closed-loop control is the "
        "sharpest benchmark for one: a wrong model shows up immediately as tracking "
        "error.",
     9.0, 2.95, 3.7, 1.3, size=14)
text(s, "Prediction accuracy is a proxy.\nControl performance is the test.",
     9.0, 4.3, 3.7, 0.9, size=16, bold=True)
text(s, "Open problems, stated plainly:", 9.0, 5.3, 3.7, 0.3, size=13, bold=True, col=RED)
bullets(s, [
    ("One session, trial-averaged; disturbance measured post-hoc", 0, False, RED),
    ("The contra predictor reconstructs — it does not yet forecast", 0, False, RED),
    ("What trajectory SHOULD activity follow?", 0, False, RED),
], 9.0, 5.65, 3.7, 1.5, size=12)
notes(s, """
THE TALK'S THESIS. Protect this slide's time.

SOURCE: ctrl_optimal_control.m on m4 (AL_0033 2025-02-26 e2), figure on disk at
paper/images/predictor_saga/ctrl_optimal_AL_0033_0226_e2.png. Plant nx=2, u_max=2.06
(PI's own usage ceiling), ref = -5, trial-averaged.

NUMBERS:
 - Steady RMSE (1-3 s): PI 1.202 -> causal optimal 0.146 (-87.8%), non-causal 0.000.
 - Full window (0-3 s): PI 1.510 / causal 0.994 / non-causal 0.927.
 - Mean laser command: PI 1.563 -> optimal 1.300.
 - Measured disturbance drifts +1.82% over the trial (waning inhibition = adaptation);
   both optima ramp u up late to fight it.

THE TWO POINTS THAT MATTER MORE THAN THE HEADLINE:
 1. BETTER TRACKING WITH LESS COMMAND. This closes the loop on slide 1's minimum-energy
    framing -- least-perturbative is not a trade against performance here. PUNCHLINE.
 2. CAUSAL (0.146) ~ CLAIRVOYANT (0.000). Nearly all the achievable gain is reachable
    CAUSALLY, so the PI gap is NOT about needing to see the future. This pre-empts the
    obvious objection that of course a non-causal optimum wins.

*** HONESTY -- SAY THESE OUT LOUD ***
 - ONE SESSION, TRIAL-AVERAGED. Not a cross-session claim.
 - THE DISTURBANCE DEFINITION FLATTERS THE OPTIMIZER. d = y_PI - H*u_PI is a post-hoc
   residual measured from the PI trials, so it ABSORBS PLANT MISMATCH -- part of the
   advantage is the optimizer repairing the plant model, not better control. The
   project's own reading rule (2026-08-12): a gap that survives under 'global' is a
   control result; one that appears only under 'residual' is the optimizer fixing the
   model.
 - THE REAL-TIME VERSION IS NOT DONE. ctrl_optimal_xsess.m with src='global' (the
   contra-predicted disturbance, which IS available in real time) has run on only
   n = 2 sessions, p = 0.5, and its exemplar is AL_0033_0212_e2 -- the session flagged
   for EXCLUSION. DO NOT show the f4_mpc_*_global panels; do not quote n=2.
 - The contra predictor RECONSTRUCTS, it does not FORECAST (pX=0), and Stage 5a shows
   contra adds no temporal information beyond ipsi's own past (0.835 vs 0.863 at the
   86 ms horizon). The forecast must come from the NETWORK DYNAMICS (slide 9).

THE SENTENCE: "on one session, with the disturbance measured post-hoc, the optimal
controller more than halves the tracking error using less light -- and we are running
the version that uses the real-time-available disturbance now."

AIRCRAFT CALLBACK (analogy's SECOND AND LAST appearance -- do not reuse the rocket):
gust-load alleviation uses forward-looking sensing to measure the gust BEFORE it hits
the wing and feeds it forward. Preview control generalised to a PREDICTED disturbance.

ASK FROM THE STAGE: you want forward models good enough to close the loop on. TIM KIM
is an organiser (latent-space widefield model, already on TASKS); GUILLAUME BELLEC is
speaking. This is the moment. Then acknowledgments.
""")

# ---------------------------------------------------------------- BACKUP
s = slide("BACKUP — Is that controller work, or state-dependent actuator gain?")
text(s, "Nick's central objection. A control-literate room lands here fast.",
     0.6, 1.5, 12.1, 0.4, size=15, italic=True)
bullets(s, [
    "The concern: the pre-stim baseline offset in the disturbance-rejection scatter "
    "may reflect trial-by-trial LASER GAIN variability — state-dependent actuator "
    "scaling — rather than controller work",
    ("Two tests, designed and NOT YET RUN:", 0, True),
    ("correlate each trial's pre-stim baseline offset against pre-stim delta power "
     "and motion", 1),
    ("regress stim-period residual amplitude ~ pre-stim delta × condition (OL/CL); "
     "a significant INTERACTION means actuator-gain variability, not controller "
     "failure", 1),
    "Same identification problem already flagged for the Fig-4 Stage-5 pair — the "
    "interaction term is the fix both need",
], 0.6, 2.05, 12.1, 3.5, size=14)
text(s, "Honest answer: designed, not yet run.", 0.6, 5.7, 12.1, 0.4, size=15,
     bold=True, col=RED)
notes(s, "Do not bluff this one. 'Designed, not yet run' is far stronger than "
         "improvising a defence. It is a real open question and Nick raised it first.")

s = slide("BACKUP — Hemodynamics")
bullets(s, [
    "Single-wavelength ΔF/F conflates neural signal with blood-volume / arousal "
    "components",
    "7 of 13 controller sessions have hemo-corrected data; 5 fall back to uncorrected "
    "+ detrendAndFilt",
    "The split is carried as a covariate — do not pool blind",
    "This is the single biggest anticipated referee objection, and it is tracked, "
    "not ignored",
], 0.6, 1.7, 12.1, 3.0, size=15)
notes(s, "Know the 7/13 number cold. This is the T1 objection in the project's own "
         "red-team list, so owning it reads as rigour.")

s = slide("BACKUP — Why PI and not MPC?")
bullets(s, [
    "Latency budget: ~14 ms minimum loop, ~160 ms plant lag",
    "MPC needs a forward model good enough to trust over the horizon — that is "
    "exactly what slide 8 is building",
    "One-sided actuation makes the optimisation a constrained QP, not LQR",
    "PI is the honest baseline: it works, it transfers across mice, and it sets the "
    "bar a model-based controller has to beat",
], 0.6, 1.7, 12.1, 3.0, size=15)
notes(s, "Optional 4th backup: gain-grid cost surfaces (T-A/T-B) + autotune "
         "convergence (T-C) if anyone asks how the gains were chosen.")

out = os.path.join(HERE, "neuroai_seattle_2026.pptx")
prs.save(out)
print("wrote", out)
print(os.path.getsize(out) // (1024 * 1024), "MB |", len(prs.slides._sldIdLst), "slides")
