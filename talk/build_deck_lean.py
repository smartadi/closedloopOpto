"""LEAN build of the NeuroAI Seattle 2026 deck.

Same slides, same order, same figures as build_deck.py -- but the slides carry only
PUNCHLINES. Every explanatory sentence lives in the speaker notes (loaded verbatim
from notes.json, which was extracted from the verbose deck, plus the prose stripped
off each slide here). Output: neuroai_seattle_2026_lean.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import json
import os

Image.MAX_IMAGE_PIXELS = None
HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
ROOT = os.path.dirname(HERE)
NOTES = json.load(open(os.path.join(HERE, "notes.json"), encoding="utf-8"))

W, H = 13.333, 7.5
BLACK = RGBColor(0x11, 0x11, 0x11)
GREY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xB0, 0x00, 0x00)
BODY = HEAD = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


def slide(title=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(10.4), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        r = tf.paragraphs[0].add_run()
        r.text = title
        r.font.size, r.font.bold, r.font.name = Pt(28), True, HEAD
        r.font.color.rgb = BLACK
    if sub:
        tb = s.shapes.add_textbox(Inches(11.1), Inches(0.38), Inches(1.65), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = sub
        r.font.size, r.font.name, r.font.italic = Pt(12), BODY, True
        r.font.color.rgb = GREY
    return s


def bullets(s, items, x, y, w, h, size=18, gap=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl, bold, col = (list(it) + [0, False, None])[:4]
        else:
            txt, lvl, bold, col = it, 0, False, None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.name = BODY
        r.font.bold = bool(bold)
        r.font.color.rgb = col or (BLACK if lvl == 0 else GREY)
    return tb


def text(s, txt, x, y, w, h, size=18, bold=False, col=None, italic=False, align=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        if align:
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), BODY, bold, italic
        r.font.color.rgb = col or BLACK
    return tb


def big(s, txt, x, y, w, size=30, col=None, align=PP_ALIGN.CENTER):
    """Number callout."""
    return text(s, txt, x, y, w, 0.7, size=size, bold=True, col=col, align=align)


def pic(s, name, x, y, boxw, boxh, cap=None, top=False):
    path = os.path.join(IMG, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = boxw, boxw / ar
    if h > boxh:
        h, w = boxh, boxh * ar
    px = x + (boxw - w) / 2
    py = y if top else y + (boxh - h) / 2
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))
    if cap:
        text(s, cap, x, py + h + 0.08, boxw, 0.3, size=12, col=GREY,
             align=PP_ALIGN.CENTER)
    return (px, py, w, h)


def caption(s, txt, x, y, w):
    text(s, txt, x, y, w, 0.3, size=12, col=GREY, align=PP_ALIGN.CENTER)


def notes(s, key, extra=""):
    """Speaker notes = the verbose deck's notes + whatever prose left the slide."""
    body = NOTES.get(key, "")
    if extra:
        body = (body + "\n\n--- MOVED OFF THE SLIDE IN THE LEAN CUT ---\n"
                + extra.strip())
    s.notes_slide.notes_text_frame.text = body.strip()


def movie(s, fname, poster, x, y, w, h):
    try:
        s.shapes.add_movie(os.path.join(ROOT, "presentation", fname),
                           Inches(x), Inches(y), Inches(w), Inches(h),
                           poster_frame_image=os.path.join(IMG, poster),
                           mime_type="video/mp4")
    except Exception as e:
        pic(s, poster, x, y, w, h)
        print("  ! movie embed failed for", fname, "->", e)


# ================================================================ 0 TITLE
s = slide()
text(s, "Closed-loop optogenetic control of\nmesoscale cortical population activity",
     0.9, 2.2, W - 1.8, 2.0, size=36, bold=True)
text(s, "Aditya  ·  NeuroAI in Seattle 2026  ·  Allen Institute  ·  20–21 August",
     0.9, 4.5, W - 1.8, 0.5, size=17, col=GREY)
notes(s, "title")

# ================================================================ 1 OPENING
s = slide("We can measure it, and we can act on it. So control it.", "1:15")
bullets(s, [
    ("Measure — widefield, cortex-wide, 35 Hz", 0, True),
    ("Act — spatially targeted optogenetics", 0, True),
    ("Identify the plant between them", 0, True),
], 0.7, 1.6, 6.0, 1.8, size=20)

text(s, "What it unlocks", 0.7, 3.6, 6.0, 0.35, size=17, bold=True, col=GREY)
bullets(s, [
    "Brain state as the feedback signal",
    "Optimization, not hand-tuning",
    "Minimum energy = least perturbative",
    "Model-driven, not tuned",
], 0.7, 4.0, 6.0, 2.0, size=18)

text(s, "Like landing a rocket", 7.4, 1.6, 5.3, 0.35, size=17, bold=True, col=GREY)
bullets(s, [
    "One-sided thrust:  0 ≤ u ≤ u_max",
    "Minimum fuel, not any landing",
    ("⇒ constrained QP, not LQR", 0, True),
], 7.4, 2.0, 5.3, 1.5, size=18)
text(s, "…but far more variable.", 7.4, 3.7, 5.3, 0.4, size=18, italic=True)

text(s, "Not “we delivered 2 mW”.\n“We held it at −5 % ΔF/F.”",
     7.4, 4.5, 5.3, 1.2, size=22, bold=True)
text(s, "The perturbation becomes a controlled variable.",
     7.4, 5.9, 5.3, 0.5, size=17, italic=True, col=GREY)

text(s, "plant → loop → residual → model",
     0.7, 6.6, 12.0, 0.4, size=17, italic=True, col=GREY)
notes(s, "opening", """
PRIOR WORK -- SAY IT, DO NOT SLIDE IT. Closed-loop optogenetics is not new:
Newman 2015 (the optoclamp, also a PI controller -- our direct ancestor), Bolus
2018/2021 (model-based, then LQR in awake mice), Kanta 2019, Clancy 2014, Grosenick
2015. DO NOT FRAME OPEN-LOOP AS FLAWED -- Svoboda is an organiser and in the room.
THE GAP, IN ONE SENTENCE: one thread has the LOOP but not the SCALE (single-neuron or
LFP, ephys feedback); the other has the THEORY but no ACTUATOR (controllability on
structural connectomes, inputs hypothetical). Widefield + targeted opto gives loop,
scale and theory at once.
""")

# ================================================================ 2 THE LOOP
s = slide("The loop", "0:45")
pic(s, "wfpath.png", 0.7, 1.5, 3.0, 4.3, top=True)
pic(s, "svd_frame.png", 4.1, 1.5, 3.0, 4.3, top=True)
pic(s, "arch.png", 7.4, 1.5, 5.3, 2.8, top=True)
big(s, "35 Hz", 7.4, 4.6, 1.7, size=26)
big(s, "14 ms", 9.2, 4.6, 1.7, size=26)
big(s, "−5 % ΔF/F", 11.0, 4.6, 1.9, size=26)
caption(s, "control rate", 7.4, 5.25, 1.7)
caption(s, "loop latency", 9.2, 5.25, 1.7)
caption(s, "reference", 11.0, 5.25, 1.9)
notes(s, "loop", """
Widefield SVD readout at 35 Hz -> kernel-mean dF/F at a target pixel -> PI controller
-> 638 / 594 nm laser. Plant lag ~160 ms -- comes back on the preview slide, where the
143 ms lookahead is matched to it BY DESIGN.
""")

# ================================================================ 3 PLANT
s = slide("The plant is approximately linear — and it transfers", "1:00")
pic(s, "imp_single.png", 0.6, 1.5, 3.9, 3.2, top=True)
pic(s, "imp_response.png", 4.7, 1.5, 3.9, 3.2, top=True)
pic(s, "tf_shape.png", 8.8, 1.5, 3.9, 3.2, top=True)
text(s, "A low-order LTI model explains the trial-averaged response.",
     0.6, 5.0, 12.1, 0.5, size=22, bold=True)
bullets(s, [
    "Similar time scales across sessions",
    "Similar shapes — not strictly the same",
], 0.6, 5.6, 12.1, 1.2, size=18)
notes(s, "plant", """
User's own wording from the verbose cut: "Input-output characterization per session
reveals that LTI systems explain trial-averaged dynamics; TF is roughly consistent
across sessions -- similar time scales, similar shapes, not strictly the same though."
SIGNPOST FORWARD: "hold onto this -- MPC needs a forward model, and this is it."
""")

# ================================================================ 4 CONTROLLER
s = slide("Feedforward inversion, then PI on top", "0:45")
text(s, "u = Kr·e  +  Kp·e  +  Ki·Σ e Δt", 0.7, 1.7, 6.2, 0.6, size=26, bold=True)
text(s, "e = y_ref − y", 0.7, 2.35, 6.2, 0.4, size=18, col=GREY)

text(s, "Open loop = the same law,\nKp = Ki = 0.", 0.7, 3.1, 6.2, 1.0,
     size=24, bold=True)
text(s, "Feedback added to a calibrated feedforward.", 0.7, 4.2, 6.2, 0.4,
     size=17, italic=True, col=GREY)

bullets(s, [
    "0 ≤ u ≤ 5 V — the one-sided bound, in hardware",
    ("The integrator is a STATE:  ξ = Σ e Δt", 0, True),
], 0.7, 5.0, 6.2, 1.4, size=18)

pic(s, "tune_grid.png", 7.3, 1.6, 2.8, 2.6, top=True)
pic(s, "tune_auto_cost.png", 10.2, 1.6, 2.5, 2.4, top=True)
caption(s, "gain grid J(Kp,Ki)", 7.3, 4.25, 2.8)
caption(s, "online auto-tune", 10.2, 4.25, 2.5)
text(s, "Tuned per session, model-free.", 7.3, 4.9, 5.4, 0.4, size=20, bold=True)
text(s, "No model needed to tune PI.\nA model is what it takes to beat it.",
     7.3, 5.5, 5.4, 1.0, size=18, italic=True)
notes(s, "ctrl", """
Kr = U Y-dagger, least squares on the impulse calibration (Methods eq:kr).
Anti-windup clamp on the accumulator while the actuator saturates.
0-5 V = 36 mW max.
Grid basin approx (0.05, 0.10); the online zeroth-order search settles at
(0.068, 0.064) -- SAME MOUSE, AL_0033. Cost J = mean ||y - y_ref||_2 over t = 0-3 s,
ref = -5 %dF/F. Auto-tune cost 16.6 -> 12.3.
""")

# ================================================================ 5 CL VIDEO
s = slide("Closed loop holds the set-point", "1:45  ·  VIDEO")
movie(s, "controller_demo_m13_full_hd.mp4", "poster_ctrl.png", 0.6, 1.6, 7.6, 4.28)
pic(s, "pooled_rmse.png", 8.7, 1.7, 4.0, 3.6)
big(s, "14 / 15 sessions", 8.7, 5.5, 4.0, size=26)
caption(s, "CL beats OL  ·  signrank p = 1.2 × 10⁻⁴", 8.7, 6.2, 4.0)
notes(s, "clvideo", """
VARIANCE + SPREAD ADDED 2026-08-18 at user request.
 - Variance: OL/CL across-trial variance ratio by window (variance_mse.m fig_Fr,
   panel 3I). Pre 0.95, STIM 1.85 ***, Post 0.98. The effect is confined to the stim
   window -- pre and post sit on 1 -- so it is not a baseline or a session-level shift.
 - Spread: per-session trial-RMSE half-violins, OL vs CL (variance_mse.m fig_G,
   panel 3H, 13 sessions). The point is that CL narrows the DISTRIBUTION, not just the
   median -- the bad trials are what disappear.
 - Medians: per-session median OL vs CL RMSE, 15 sessions (3K). Pooled OL 3.11 /
   CL 2.44 %dF/F, RMSE window t = 0 -> +3 s.
 - OL stim alone already reduces variance a little (post-onset slope -0.57 +/- 0.24
   vs -0.04 +/- 0.14 pre, n=13). CL adds a further, more consistent reduction. Do not
   claim stimulation has no effect on variance -- that older claim was WRONG.
 - Video = m13 (AL_0039 2025-04-20), step, 100 trials, 0:39. ~30% CL-vs-OL gap.
CAVEAT ON THE VIOLIN PANEL: it is drawn red = OL / GREEN = CL, which is not the
project's locked colour scheme (OL red, CL BLUE). Re-export before the final deck, or
say "red is open loop, green is closed loop" out loud once.
""")

# ================================================================ 6 DRIVERS
s = slide("What drives the residual error", "1:00")
pic(s, "f4_decomp.png", 0.5, 1.5, 4.1, 3.0, top=True)
pic(s, "f4_initdev.png", 4.8, 1.5, 3.9, 3.0, top=True)
pic(s, "f4_hi24.png", 8.9, 1.5, 3.9, 3.0, top=True)
caption(s, "unique R² of trial RMSE  ·  613 CL trials / 11 sessions", 0.5, 4.65, 12.3)

text(s, "Early: the transient.", 0.7, 5.3, 3.7, 0.4, size=20, bold=True)
big(s, "R² 0.38 → 0.03", 0.7, 5.75, 3.7, size=22, align=PP_ALIGN.LEFT)

text(s, "Motion: nothing.", 5.0, 5.3, 3.3, 0.4, size=20, bold=True)
big(s, "R² 0.001", 5.0, 5.75, 3.3, size=22, align=PP_ALIGN.LEFT)

text(s, "Late: 2–4 Hz state.", 8.9, 5.3, 3.9, 0.4, size=20, bold=True)
big(s, "93 % of explained", 8.9, 5.75, 3.9, size=22, align=PP_ALIGN.LEFT)
notes(s, "drivers", """
Not a tuning problem, and not movement. It is one specific band of ongoing cortical
activity -- which is what the next slide is about.
Unique partial R2, 2000-bootstrap 95% CI: initial deviation 0.381 [0.299 0.458] early
-> 0.029 [0.004 0.077] late (13x collapse); motion 0.001 -> 0.001; relative 2-4 Hz
0.089 -> 0.122. Full model R2 0.407 early / 0.132 late. Relative 2-4 Hz = 22% of
explained variance early, 93% LATE.
ROBUSTNESS: motion gate off (852 trials / 15 sessions) gives R2 0.389/0.123 -- nothing
moves. That licenses "motion does not matter" over "motion was not measured".
""")

# ================================================================ 7 WATERBED
s = slide("The loop rejects DC — and makes 2–4 Hz worse", "1:15")
pic(s, "f4_slope_motion.png", 0.5, 1.5, 3.5, 2.6, top=True)
pic(s, "f4_slope_delta24.png", 4.2, 1.5, 3.5, 2.6, top=True)
pic(s, "f4_slope_delta12.png", 7.9, 1.5, 3.5, 2.6, top=True)
caption(s, "motion  +0.69 → −0.02   p = 0.002", 0.5, 4.3, 3.5)
caption(s, "2–4 Hz  −0.51 → +0.81   p = 0.013", 4.2, 4.3, 3.5)
caption(s, "1–2 Hz  null control   p = 0.76", 7.9, 4.3, 3.5)
caption(s, "per-session slope of trial RMSE on each state factor, OL vs CL — "
           "15 sessions, four mice", 0.5, 4.65, 10.9)

text(s, "The integrator is a model of a constant — so DC goes away.",
     0.6, 5.3, 12.1, 0.4, size=20, bold=True)
text(s, "There is no model of 2–4 Hz — so the slope FLIPS SIGN.",
     0.6, 5.85, 12.1, 0.4, size=20, bold=True)
text(s, "Not a null result. That is the waterbed.",
     0.6, 6.45, 12.1, 0.5, size=22, bold=True, italic=True, col=RED)
notes(s, "reject_what", """
INTERNAL MODEL PRINCIPLE, stated aloud here for the first time: to reject a
disturbance, the loop must contain a model of its dynamics.
A controller that merely LACKED a model of that band would leave the OL slope alone
(S ~ 1: no rejection, no harm). A FLIP means feedback is actively making 2-4 Hz worse
-- it is paying for its low-frequency suppression. Bode's sensitivity integral.
1-2 Hz is the null control, so the dissociation is band-specific, not generic.
PANEL HYGIENE: the 1-2 Hz null MUST be the per-session paired slope
(factor_slope_delta12), NOT claim2_delta_lo12_late -- in the pooled-quartile form the
1-2 Hz series rises almost identically to 2-4 Hz and argues AGAINST the dissociation.
""")

# ================================================================ 8 MAP
s = slide("The disturbance, measured: a stim-blind contra → ipsi map", "0:45")
pic(s, "kernel_paint.png", 0.5, 1.4, 6.3, 5.5, top=True)
caption(s, "six ipsi target sites (ringed) and their contra kernels", 0.5, 7.0, 6.3)

text(s, "Topographic — not one diffuse global signal.",
     7.3, 1.7, 5.4, 0.8, size=22, bold=True)

bullets(s, [
    "Fit on spontaneous, laser-off frames only",
    "Contra pixels the laser cannot reach ⇒ stim-blind",
    ("Global = predicted ipsi, as if no laser", 0, True),
    ("Local = Actual − Global", 0, True),
], 7.3, 2.9, 5.4, 2.4, size=18)

big(s, "R² = 0.87", 7.3, 5.3, 5.4, size=30, align=PP_ALIGN.LEFT)
caption(s, "held-out contra → ipsi (m4); 0.978 on Ye et al. 2023's own data",
        7.3, 5.95, 5.4)
text(s, "Reject a disturbance and the disturbance\nis the rest of the brain.",
     7.3, 6.35, 5.4, 0.9, size=18, bold=True, italic=True)
notes(s, "map", """
Each ipsi site is predicted by a distinct, spatially coherent territory of
contralateral cortex. The dots are the sparse predictor grid, coloured by which ipsi
target that contra pixel best predicts; opacity is proportional to |weight|.
The fit is deployed during trials to give Global -- the counterfactual ipsi trace
"as if no laser" -- and Local = Actual - Global.
""")

# ================================================================ 9 IMP
s = slide("Feedback rejects it. The floor is a missing model.", "1:30")
pic(s, "rej_demo_ol.png", 0.6, 1.5, 3.3, 2.1, cap="OPEN LOOP", top=True)
pic(s, "rej_demo_cl.png", 0.6, 3.9, 3.3, 2.1, cap="CLOSED LOOP", top=True)
pic(s, "rej_paired.png", 4.3, 1.9, 3.1, 2.7, top=True)
caption(s, "energy ratio  ER = ‖A − ref‖² / ‖G − ref‖²    1 = no work done",
        4.0, 4.8, 3.9)
big(s, "0.291 → 0.203", 4.0, 5.3, 3.9, size=26)
caption(s, "13 of 13 sessions  ·  p = 2.4 × 10⁻⁴", 4.0, 6.0, 3.9)

text(s, "PI's integrator is an internal model\nof a CONSTANT.",
     8.0, 1.7, 4.7, 1.0, size=22, bold=True)
text(s, "Ongoing cortical activity is not constant.",
     8.0, 2.9, 4.7, 0.5, size=20)
text(s, "So it stops at ER ≈ 0.20.", 8.0, 3.5, 4.7, 0.5, size=20)

text(s, "The residual is not a tuning failure.\nIt is a missing model.",
     8.0, 4.6, 4.7, 1.1, size=24, bold=True, italic=True, col=RED)
caption(s, "panels are UNGATED (no R² floor, n = 13)", 0.6, 6.2, 3.3)
notes(s, "imp", """
Per-trial transmission beta (slope of Actual on Global): m4 OL 0.97 -> CL 0.56.
Open-loop the disturbance passes through untouched; feedback nearly halves it.
Rejection is state-INVARIANT -- a constant OL-CL offset across motion and delta
quartiles. State-dependence returned NULL; do not imply a state effect.
Per-session rank-sum significant in 11/13; the two misses are the two new mice with
weak predictors, so their inclusion is CONSERVATIVE.
==> MPC is not an add-on. It is what the internal model principle prescribes.
""")

# ================================================================ 10 SINE
s = slide("A moving reference: feedback fixes error, preview fixes lag",
          "1:15  ·  VIDEO")
movie(s, "ff_demo_0721_hd.mp4", "poster_ff.png", 0.6, 1.6, 6.6, 3.71)
pic(s, "sine_rmse.png", 7.6, 1.6, 2.4, 3.1, top=True)
pic(s, "sine_phase.png", 10.3, 1.6, 2.4, 3.1, top=True)
caption(s, "total RMSE", 7.6, 4.8, 2.4)
caption(s, "1 Hz phase lag", 10.3, 4.8, 2.4)

text(s, "Error ← FEEDBACK", 0.7, 5.6, 4.0, 0.5, size=22, bold=True)
text(s, "Lag ← PREVIEW", 5.0, 5.6, 4.0, 0.5, size=22, bold=True)
big(s, "60° → 5°", 9.4, 5.5, 3.3, size=26)
caption(s, "143 ms lookahead ≈ the 160 ms plant lag", 9.4, 6.15, 3.3)
text(s, "Preliminary: 3 sessions, one mouse, one hemisphere.",
     0.7, 6.6, 8.0, 0.4, size=16, col=RED, italic=True)
notes(s, "sine", """
CL < OL in all 3 sessions, pooled p = 1.5e-4.
DO NOT SAY: that preview reduces error (OL vs OL+preview N.S., p = 0.49); that
CL+preview beats CL (p = 0.093 in CL's favour); that Friedman across sessions is
significant (n = 3, p = 0.060, underpowered).
The supportable sentence: "CL+preview is the only mode that is both significantly
below open-loop on error AND lag-free."
WHY IT IS HERE: the only place in the talk where LOOKAHEAD works on real data --
an existence proof for anticipation on a KNOWN reference, which the close then
generalises to a PREDICTED disturbance.
""")

# ================================================================ 11 CLOSE
s = slide("Give it a model and it does better — with less light", "1:15")
pic(s, "opt_tracking.png", 0.6, 1.5, 4.0, 3.3, top=True)
pic(s, "opt_command.png", 4.8, 1.5, 3.9, 3.3, top=True)
caption(s, "m4, trial-averaged — PI overshoots and oscillates; "
           "the optimal controller sits on the reference", 0.6, 4.9, 8.1)

big(s, "RMSE  1.20 → 0.15", 0.6, 5.4, 4.0, size=24)
big(s, "light  1.56 → 1.30", 4.8, 5.4, 3.9, size=24)
caption(s, "PI → optimal", 0.6, 6.1, 4.0)
caption(s, "better tracking, less light", 4.8, 6.1, 3.9)
text(s, "Causal 0.146 ≈ clairvoyant 0.000 — the gap is not about seeing the future.",
     0.6, 6.55, 8.1, 0.4, size=16, col=GREY)

text(s, "Prediction accuracy is a proxy.\nControl performance is the test.",
     9.0, 1.7, 3.7, 1.2, size=24, bold=True)
text(s, "A wrong model shows up immediately as tracking error.",
     9.0, 3.1, 3.7, 0.8, size=17, italic=True, col=GREY)
text(s, "Open:", 9.0, 4.5, 3.7, 0.3, size=17, bold=True, col=RED)
bullets(s, [
    ("One session, trial-averaged", 0, False, RED),
    ("The predictor reconstructs, it does not forecast", 0, False, RED),
    ("What trajectory SHOULD activity follow?", 0, False, RED),
], 9.0, 4.85, 3.7, 2.0, size=16)
notes(s, "close", """
THE ARGUMENT: PI rejects a constant. Give the controller a model of the network's
dynamics and it rejects what the network actually does. That model is where learned
methods enter -- and closed-loop control is the sharpest benchmark for one.
Steady RMSE PI 1.20 -> optimal 0.15. Mean laser command 1.563 -> 1.300.
Least-perturbative and best-performing are NOT in tension here.
""")

# ================================================================ BACKUP
s = slide("BACKUP — Controller work, or state-dependent actuator gain?")
text(s, "Nick's central objection.", 0.7, 1.6, 12.0, 0.5, size=20, italic=True)
bullets(s, [
    "The pre-stim baseline offset may be trial-by-trial LASER GAIN, not controller work",
    "Test 1: baseline offset vs pre-stim δ power and motion",
    "Test 2: residual amplitude ~ pre-stim δ × condition — a significant INTERACTION "
    "means actuator gain",
], 0.7, 2.3, 12.0, 3.0, size=18)
text(s, "Designed, not yet run.", 0.7, 5.5, 12.0, 0.5, size=24, bold=True, col=RED)
notes(s, "bk_gain")

s = slide("BACKUP — Hemodynamics")
bullets(s, [
    "Single-wavelength ΔF/F conflates neural signal with blood volume / arousal",
    "7 of 13 controller sessions are hemo-corrected; 5 fall back to uncorrected",
    "Carried as a covariate — not pooled blind",
], 0.7, 1.9, 12.0, 3.0, size=20)
text(s, "The biggest anticipated referee objection — tracked, not ignored.",
     0.7, 5.0, 12.0, 0.5, size=20, bold=True)
notes(s, "bk_hemo")

s = slide("BACKUP — Why PI and not MPC?")
bullets(s, [
    "No per-animal forward model when the loop was built",
    "PI needs no model — and the grid + auto-tune agree without one",
    "One-sided actuation ⇒ constrained QP, not LQR",
    "PI is the bar a model-based controller has to beat",
], 0.7, 1.9, 12.0, 3.4, size=20)
notes(s, "bk_mpc")

out = os.path.join(HERE, "neuroai_seattle_2026_lean.pptx")
prs.save(out)
print("wrote", out)
print(os.path.getsize(out) // (1024 * 1024), "MB |", len(prs.slides._sldIdLst), "slides")
